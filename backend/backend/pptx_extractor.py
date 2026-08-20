"""PPTX slide text and image extraction.

Public API
----------
extract_pptx_slides(pptx_path) -> list[dict]
    Returns per-slide records:
        {
            "slide_num":    int,
            "text":         str,
            "image_base64": str | None   # base64 JPEG or None
        }
"""

from __future__ import annotations

import base64
import io

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

# Maximum slides for which images are generated.
MAX_SLIDE_IMAGES = 30

# Target width (pixels) for all slide preview images.
_TARGET_W = 800

# JPEG quality (0-95).  72 gives good quality at ~100-200 KB per slide.
_JPEG_QUALITY = 72


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------

def _pil_to_jpeg_b64(img: Image.Image) -> str | None:
    """Resize to <= _TARGET_W, encode as JPEG, return base64 string or None."""
    try:
        if img.width > _TARGET_W:
            ratio = _TARGET_W / img.width
            new_h = max(1, int(img.height * ratio))
            img = img.resize((_TARGET_W, new_h), Image.LANCZOS)
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")
        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=_JPEG_QUALITY, optimize=True)
        return base64.b64encode(buf.getvalue()).decode("utf-8")
    except Exception:
        return None


def _load_font(size: int):
    """Return a PIL ImageFont at the given size, or the default bitmap font."""
    candidates = (
        "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/ubuntu/Ubuntu-R.ttf",
        "/usr/share/fonts/truetype/freefont/FreeSans.ttf",
        "C:/Windows/Fonts/arial.ttf",
        "C:/Windows/Fonts/Arial.ttf",
    )
    for path in candidates:
        try:
            return ImageFont.truetype(path, size)
        except Exception:
            continue
    try:
        return ImageFont.load_default()
    except Exception:
        return None


def _shape_bounds(shape, prs_w: int, prs_h: int, tgt_w: int, tgt_h: int):
    """Convert shape EMU coordinates to pixel rect clamped to canvas.

    Returns (left, top, width, height) or None on error.
    """
    try:
        left = int((shape.left or 0) / prs_w * tgt_w)
        top = int((shape.top or 0) / prs_h * tgt_h)
        w = max(1, int((shape.width or 0) / prs_w * tgt_w))
        h = max(1, int((shape.height or 0) / prs_h * tgt_h))
        left = max(0, min(left, tgt_w - 1))
        top = max(0, min(top, tgt_h - 1))
        w = min(w, tgt_w - left)
        h = min(h, tgt_h - top)
        if w < 1 or h < 1:
            return None
        return left, top, w, h
    except Exception:
        return None


def _is_blank_canvas(img: Image.Image) -> bool:
    """Return True when the canvas contains only white pixels."""
    try:
        for channel_range in img.getextrema():
            lo, hi = channel_range if isinstance(channel_range, tuple) else (channel_range, channel_range)
            if lo != 255 or hi != 255:
                return False
        return True
    except Exception:
        return False


def _render_slide_composite(slide, prs_w: int, prs_h: int) -> str | None:
    """Composite-render a slide: paste pictures and draw text shapes.

    Strategy
    --------
    1. White canvas sized proportionally to the presentation dimensions.
    2. Iterate shapes in z-order (pptx shape list order).
    3. PICTURE shapes  -> paste the image blob at the proportional position.
    4. Text-frame shapes -> draw the paragraph text using PIL.
    5. If the final canvas is pure white (no visible content rendered),
       return None so the fallback path runs.

    All per-shape errors are silently ignored so one bad shape never
    aborts rendering for the whole slide.
    """
    if not prs_w or not prs_h:
        return None
    try:
        tgt_h = max(1, int(prs_h * _TARGET_W / prs_w))
        canvas = Image.new("RGB", (_TARGET_W, tgt_h), (255, 255, 255))
        draw = ImageDraw.Draw(canvas)
        font_title = _load_font(22)
        font_body = _load_font(14)
        has_content = False

        for shape in slide.shapes:
            try:
                bounds = _shape_bounds(shape, prs_w, prs_h, _TARGET_W, tgt_h)
                if bounds is None:
                    continue
                left, top, w, h = bounds

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    pic = Image.open(io.BytesIO(shape.image.blob))
                    if pic.mode not in ("RGB", "L"):
                        pic = pic.convert("RGB")
                    pic = pic.resize((w, h), Image.LANCZOS)
                    canvas.paste(pic, (left, top))
                    has_content = True

                elif shape.has_text_frame:
                    lines = [
                        para.text.strip()
                        for para in shape.text_frame.paragraphs
                        if para.text.strip()
                    ]
                    if lines:
                        font = font_title if top < tgt_h * 0.25 else font_body
                        draw.text(
                            (left + 4, top + 4),
                            "\n".join(lines),
                            fill=(20, 20, 20),
                            font=font,
                        )
                        has_content = True

            except Exception:
                continue

        if not has_content or _is_blank_canvas(canvas):
            return None

        return _pil_to_jpeg_b64(canvas)
    except Exception:
        return None


def _extract_largest_picture(slide) -> str | None:
    """Fallback: return the largest embedded PICTURE from a slide."""
    try:
        best_area = 0
        best_blob: bytes | None = None
        for shape in slide.shapes:
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    area = (shape.width or 0) * (shape.height or 0)
                    if area > best_area:
                        best_area = area
                        best_blob = shape.image.blob
            except Exception:
                continue
        if best_blob is None:
            return None
        img = Image.open(io.BytesIO(best_blob))
        return _pil_to_jpeg_b64(img)
    except Exception:
        return None


def _extract_slide_text(slide) -> str:
    """Extract all readable text from a slide's shapes."""
    parts: list[str] = []
    for shape in slide.shapes:
        try:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        parts.append(t)
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                for row in shape.table.rows:
                    for cell in row.cells:
                        t = cell.text.strip()
                        if t:
                            parts.append(t)
        except Exception:
            continue
    return "\n".join(parts)


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def extract_pptx_slides(pptx_path: str) -> list[dict]:
    """Extract per-slide text and a composite preview image from a PPTX file.

    Image generation
    ----------------
    For each slide (up to MAX_SLIDE_IMAGES):
      1. Attempt a composite PIL render (pictures + text at correct positions).
      2. Fall back to the largest embedded PICTURE shape if composite is blank.
      3. If nothing is usable, image_base64 is None.

    Text extraction covers ALL slides regardless of the image limit.

    All errors are isolated per shape and per slide.

    Returns
    -------
    list of dicts:
        {
            "slide_num":    int,
            "text":         str,
            "image_base64": str | None   # base64 JPEG
        }
    """
    prs = Presentation(pptx_path)
    prs_w = int(prs.slide_width or 0)
    prs_h = int(prs.slide_height or 0)
    results: list[dict] = []

    for i, slide in enumerate(prs.slides, start=1):
        slide_text = _extract_slide_text(slide)

        image_b64: str | None = None
        if i <= MAX_SLIDE_IMAGES:
            image_b64 = _render_slide_composite(slide, prs_w, prs_h)
            if image_b64 is None:
                image_b64 = _extract_largest_picture(slide)

        results.append({
            "slide_num": i,
            "text": slide_text,
            "image_base64": image_b64,
        })

    return results
